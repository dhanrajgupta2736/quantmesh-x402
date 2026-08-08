import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Dark Color Palette
    NAVY_BG = RGBColor(9, 13, 22)
    CARD_BG = RGBColor(19, 27, 44)
    CYAN_ACCENT = RGBColor(6, 182, 212)
    PURPLE_ACCENT = RGBColor(168, 85, 247)
    GREEN_ACCENT = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(203, 213, 225)
    SLATE_MUTED = RGBColor(148, 163, 184)
    BORDER_COLOR = RGBColor(30, 41, 59)

    def set_dark_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = NAVY_BG

    def add_header(slide, title_text, subtitle_text="AlgoVerse 2026 Hackathon — PS0404"):
        # Tag
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        
        # Title
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

    def add_card(slide, left, top, width, height, title, items, color=CYAN_ACCENT):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        for item in items:
            p2 = tf.add_paragraph()
            p2.text = item
            p2.font.size = Pt(12)
            p2.font.color.rgb = LIGHT_GRAY
            p2.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s1)
    
    tb = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "QuantMesh x402"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    p2 = tf.add_paragraph()
    p2.text = "Atomic Multi-Agent Service Router on Algorand"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "Problem Statement: PS0404 — Agentic Payments (x402 Protocol)"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = PURPLE_ACCENT
    p3.space_before = Pt(18)

    p4 = tf.add_paragraph()
    p4.text = "Team: QuantMesh Innovators\nCollege: CSBS Dept., BV(DU)COE, Pune\nEvent: AlgoVerse 2026 Hackathon | Algorand Blockchain Club\nLive API: https://api.dhanrajgupta.xyz"
    p4.font.size = Pt(13)
    p4.font.color.rgb = SLATE_MUTED
    p4.space_before = Pt(25)

    if os.path.exists("assets/ai-quant-agent.png"):
        s1.shapes.add_picture("assets/ai-quant-agent.png", Inches(8.5), Inches(1.5), width=Inches(4.2))

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Motivation
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s2)
    add_header(s2, "1. Problem Statement & Motivation")

    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "❌ Subscription Fatigue",
             ["• Crypto signal APIs cost $100–$500/month.",
              "• Retail traders & bots pay high fixed costs even for 2–3 trades/mo.",
              "• Prohibitive barrier for micro autonomous AI agents."],
             CYAN_ACCENT)

    add_card(s2, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "❌ Single-AI Risk & Wasted Fees",
             ["• Single AI models hallucinate or fail under high volatility.",
              "• Standard x402 routers charge micro-fees BEFORE calling downstream workers.",
              "• User money is lost if a downstream worker crashes."],
             PURPLE_ACCENT)

    add_card(s2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "❌ Multi-Agent Friction",
             ["• Querying 4 separate AI APIs requires 4 separate transactions.",
              "• High gas fees & network latency destroy micro-payments.",
              "• Zero consumer protection for API buyers."],
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 3: Proposed Solution / Idea
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s3)
    add_header(s3, "2. Proposed Solution: QuantMesh x402 Architecture")

    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "⚡ Pay-Per-Use AI Intelligence ($0.007)",
             ["• Granular micropayments via Algorand x402 AVM scheme.",
              "• 🛡️ Pre-Execution Zero-Fee Guarantee: Router validates all 4 sub-agents FIRST. If any fail, returns HTTP 502 with $0 charged.",
              "• 🔐 Cryptographic SHA-256 Attestation: Produces verifiable signal digest anchored to transaction ID.",
              "• Fuses NLP Sentiment + Whale Flow + TA Indicators into a 0-100 composite score."],
             CYAN_ACCENT)

    add_card(s3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "🔑 Core Zero-Fee Execution Flow",
             ["1. Client sends POST request to Hono Router.",
              "2. Router executes Worker A, B, C, D in parallel.",
              "3. IF any worker fails -> HTTP 502 Sub-Agent Failed ($0 charged to user).",
              "4. IF all workers pass -> Issue HTTP 402 Challenge.",
              "5. Lute Wallet signs $0.007 USDC -> Algorand Testnet Settlement."],
             PURPLE_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 4: Identified Paying User & Business Model
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s4)
    add_header(s4, "3. Identified Paying User & Business Model")

    add_card(s4, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "🤖 Autonomous AI Agents",
             ["• Agentic trading scripts & DeFi bots.",
              "• Require low-latency, pay-per-use market intelligence.",
              "• Cannot manage $300/mo SaaS subscriptions."],
             CYAN_ACCENT)

    add_card(s4, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "📊 Quant Retail Traders",
             ["• Traders seeking high-confidence multi-agent signals.",
              "• Pay micro-fees only when executing active strategies ($0.007/signal).",
              "• Guaranteed zero fees on failed API responses."],
             PURPLE_ACCENT)

    add_card(s4, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "💼 Sustainable Business Model",
             ["• $0.007 USDC / Execution",
              "• $0.005 split to sub-agent model providers.",
              "• $0.002 router protocol fee.",
              "• 100% Zero-Cost on failed downstream calls."],
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 5: x402 Payment Flow (Technical Core)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s5)
    add_header(s5, "4. x402 Payment Flow (Technical Core)")

    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.2), Inches(5.0),
             "🔄 The 4-Step AVM Protocol",
             ["1. Challenge (402): Gateway validates workers, then returns HTTP 402 Payment Required.",
              "2. Sign: Lute Wallet signs Algorand ASA payment ($0.007 USDC, ASA 10458941).",
              "3. Retry: Client retries request with signed transaction payload.",
              "4. Settle: Gateway verifies tx on Algorand Indexer & releases signal receipt."],
             CYAN_ACCENT)

    if os.path.exists("assets/x402-payment-flow.png"):
        s5.shapes.add_picture("assets/x402-payment-flow.png", Inches(6.3), Inches(1.8), width=Inches(6.2))

    # -------------------------------------------------------------
    # SLIDE 6: Architecture & Tech Stack
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s6)
    add_header(s6, "5. System Architecture & Tech Stack")

    if os.path.exists("assets/architecture-diagram.png"):
        s6.shapes.add_picture("assets/architecture-diagram.png", Inches(0.8), Inches(1.6), width=Inches(11.7))

    # -------------------------------------------------------------
    # SLIDE 7: Demo Screenshots & Transaction Proof
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s7)
    add_header(s7, "6. Live Demo & Lora Explorer Proof")

    if os.path.exists("assets/dashboard-screenshot.png"):
        s7.shapes.add_picture("assets/dashboard-screenshot.png", Inches(0.8), Inches(1.8), width=Inches(5.6))

    if os.path.exists("assets/lora-txn-explorer.png"):
        s7.shapes.add_picture("assets/lora-txn-explorer.png", Inches(6.8), Inches(1.8), width=Inches(5.7))

    # -------------------------------------------------------------
    # SLIDE 7b: Lora Explorer Group Transaction Visual Flow
    # -------------------------------------------------------------
    s7b = prs.slides.add_slide(blank_layout)
    set_dark_bg(s7b)
    add_header(s7b, "6b. Lora Explorer Atomic Group Visual Flow")

    add_card(s7b, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0),
             "⚡ Atomic AVM Group Settlement",
             ["• Group Hash: bv4pwvqPv/ULbbpzAObL07qsDmcUGnvUzaBjwHL9Hl8=",
              "• Tx ID: 3V3VHVA3PEUZTKYCJUD7HNNGYK2XLLR7ROPTJ2HFQGTWGDBCHAMQ",
              "• App ID: 600011882 | Block: #66083763",
              "• All sub-agent calls execute atomically inside AVM group.",
              "• Zero Wasted Gas: If any single tx fails, whole group reverts."],
             CYAN_ACCENT)

    if os.path.exists("assets/lora-group-flow.png"):
        s7b.shapes.add_picture("assets/lora-group-flow.png", Inches(5.6), Inches(1.8), width=Inches(6.9))

    # -------------------------------------------------------------
    # SLIDE 8: Innovation & Differentiation
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s8)
    add_header(s8, "7. Innovation & Competitive Differentiation")

    add_card(s8, Inches(0.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🛡️ Zero-Fee Shield",
             ["• Pre-execution gating aborts at zero cost if sub-agents crash.",
              "• Prevents money loss."],
             CYAN_ACCENT)

    add_card(s8, Inches(3.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🤖 Multi-Agent Consensus",
             ["• Combines FinBERT NLP + CoinGecko Whale Flow + TA Engine.",
              "• Robust 0-100 composite score."],
             PURPLE_ACCENT)

    add_card(s8, Inches(6.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "⚡ Sub-3s AVM Finality",
             ["• Algorand AVM provides instant settlement.",
              "• $0.00018 gas fee per tx."],
             GREEN_ACCENT)

    add_card(s8, Inches(9.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🔐 Cryptographic Digest",
             ["• SHA-256 attestation digest anchored on-chain.",
              "• Auditably verifiable."],
             CYAN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 9: Completeness & Functionality
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s9)
    add_header(s9, "8. Completeness & Functionality (Done vs. Pending)")

    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "✅ Live & Operational (Done)",
             ["• Hono Gateway live on AWS EC2.",
              "• 4 Sub-Agents (FinBERT, Whale Flow, TA Engine, Fusion).",
              "• Lute Wallet x402 AVM payment integration.",
              "• Pre-Execution Zero-Fee Shield (502 logic verified).",
              "• SHA-256 Attestation receipt generation.",
              "• Next.js 16 Dashboard with Radial Score Arc."],
             GREEN_ACCENT)

    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "🔮 Future Roadmap (Pending)",
             ["• Algorand Mainnet Deployment (ASA 315667040).",
              "• PyTeal Box Storage Contract deployment.",
              "• Developer SDK (@quantmesh/x402-client).",
              "• Multi-asset automated rebalancing."],
             PURPLE_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 10: Real-World Impact & Scalability
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s10)
    add_header(s10, "9. Real-World Impact & Scalability")

    add_card(s10, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "🌍 Democratic Access",
             ["• Democratizes quant trading signals for micro-traders & AI agents.",
              "• 100x cost reduction vs monthly SaaS."],
             CYAN_ACCENT)

    add_card(s10, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "⚡ Mainnet Readiness",
             ["• Production ready today.",
              "• Mainnet switch requires updating 2 environment variables (Network & ASA ID)."],
             PURPLE_ACCENT)

    add_card(s10, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "🚀 Enterprise Scalability",
             ["• Hono + EC2 handles 1,000+ req/sec.",
              "• Algorand AVM handles 10,000 TPS with sub-second finality."],
             GREEN_ACCENT)

    prs.save("presentation.pptx")
    try:
        prs.save("docs/QuantMesh_x402_Hackathon_Pitch.pptx")
    except Exception as e:
        print(f"Note: docs save skipped: {e}")
    print("Saved presentation.pptx successfully!")

if __name__ == "__main__":
    build_presentation()
