import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Cyberpunk Dark / Institutional Quant
    NAVY_BG = RGBColor(11, 15, 25)
    CARD_BG = RGBColor(19, 27, 44)
    CYAN_ACCENT = RGBColor(6, 182, 212)
    PURPLE_ACCENT = RGBColor(168, 85, 247)
    GREEN_ACCENT = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(203, 213, 225)
    SLATE_MUTED = RGBColor(148, 163, 184)
    BORDER_COLOR = RGBColor(30, 41, 59)

    blank_layout = prs.slide_layouts[6]

    def set_dark_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = NAVY_BG

    def add_header(slide, title_text, category_text="AlgoVerse 2026 — PS0404"):
        # Category Tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        
        # Title Text
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

    def add_card(slide, left, top, width, height, title, content_bullets, accent_color=CYAN_ACCENT):
        # Card Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1.5)

        # Content Box inside
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent_color

        for bullet in content_bullets:
            p_b = tf.add_paragraph()
            p_b.text = bullet
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = LIGHT_GRAY
            p_b.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide1)
    
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "QuantMesh x402"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    p2 = tf.add_paragraph()
    p2.text = "Decentralized AI Micropayment Signal Router on Algorand"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Problem Statement: PS0404 — Atomic Multi-Agent Service Router (x402)"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = PURPLE_ACCENT
    p3.space_before = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "Team: QuantMesh Innovators  |  College: CSBS Dept., BV(DU)COE, Pune  |  Event: AlgoVerse 2026 Hackathon"
    p4.font.size = Pt(13)
    p4.font.color.rgb = SLATE_MUTED
    p4.space_before = Pt(30)

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Motivation
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide2)
    add_header(slide2, "Problem Statement & Real-World Motivation")

    add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "1. Subscription Fatigue", 
             ["• Crypto signal APIs cost $100–$500/month.",
              "• Traders pay high monthly fees even if they trade only 2–3 times.",
              "• High financial barrier for small automated bots & retail traders."], 
             CYAN_ACCENT)

    add_card(slide2, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "2. Single-AI Risk", 
             ["• Single AI models hallucinate or misread market volatility.",
              "• Single-indicator tools cause heavy trading drawdown losses.",
              "• Lack of multi-agent consensus validation."], 
             PURPLE_ACCENT)

    add_card(slide2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "3. Payment Friction", 
             ["• Querying 4 separate AI APIs requires 4 transactions.",
              "• High gas fees & friction on legacy chains.",
              "• No zero-fee protection if an API fails."], 
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 3: Proposed Solution / Idea
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide3)
    add_header(slide3, "Proposed Solution: QuantMesh x402 Architecture")

    add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "⚡ Pay-Per-Use Micropayments ($0.007)", 
             ["• Zero monthly subscription commitments.",
              "• Users pay flat $0.007 USDC/ALGO per fused signal request.",
              "• Pre-Execution Gating (Zero-Fee Guarantee): If any sub-agent fails, HTTP 502 returns & $0 is charged."], 
             CYAN_ACCENT)

    add_card(slide3, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "🤖 Multi-Agent Consensus Mesh", 
             ["• Combines 4 specialized AI sub-agents in parallel:",
              "  1. Worker A: FinBERT Financial Sentiment NLP",
              "  2. Worker B: CoinGecko On-Chain Whale Net Flow",
              "  3. Worker C: Technical Indicators (RSI, SMA, MACD)",
              "  4. Worker D: Dynamic Weighted Consensus Fusion Engine"], 
             PURPLE_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 4: Identified Paying User & Business Model
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide4)
    add_header(slide4, "Paying Users & Sustainable Unit Economics")

    add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "🎯 Target Customers & Paying Users", 
             ["• Algorithmic Crypto Traders & Telegram Signal Bots.",
              "• Web3 DeFi Protocols & DEX Interface Aggregators.",
              "• Autonomous AI Agents needing verified market data.",
              "• Pay-Per-Use model saves users 99.9% vs $300/mo subscriptions."], 
             CYAN_ACCENT)

    add_card(slide4, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "💰 Unit Economics & Protocol Margin", 
             ["• User Payment: $0.0070 USDC per execution.",
              "• Worker Payouts: $0.0060 total ($0.002 A, $0.002 B, $0.001 C, $0.001 D).",
              "• Router Net Profit: $0.0010 USDC per call (14.3% Margin).",
              "• At 1M calls/day = $1,000/day ($365k/year) in automated profit."], 
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 5: x402 Payment Flow (Technical Core)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide5)
    add_header(slide5, "x402 Technical Protocol Flow (Challenge → Sign → Retry → Settle)")

    steps = [
        ("Step 1: Probe", "Client sends POST /orchestrate without payment proof.", CYAN_ACCENT),
        ("Step 2: Challenge (402)", "Server pre-executes agents & returns HTTP 402 + payment headers.", PURPLE_ACCENT),
        ("Step 3: Sign (Lute)", "Client reads 402 headers & signs 5-txn group in Lute Wallet.", GREEN_ACCENT),
        ("Step 4: Settle & Result", "Client retries with payment proof; GoPlausible verifies & settles in 1 block.", CYAN_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(steps):
        add_card(slide5, Inches(0.8 + i * 2.95), Inches(1.8), Inches(2.8), Inches(5.0), title, [desc], color)

    # -------------------------------------------------------------
    # SLIDE 6: Architecture & Tech Stack
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide6)
    add_header(slide6, "Full-Stack Architecture & Technology Stack")

    add_card(slide6, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "🌐 Frontend Layer", 
             ["• Next.js 16 (App Router)",
              "• Tailwind CSS Glassmorphism",
              "• @txnlab/use-wallet-react",
              "• Lute / Pera Testnet Signer"], 
             CYAN_ACCENT)

    add_card(slide6, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "⚙️ Orchestrator Backend", 
             ["• Hono.js High-Speed Gateway",
              "• GoPlausible Facilitator Client",
              "• Algorand Indexer Verification",
              "• AWS EC2 Live Host"], 
             PURPLE_ACCENT)

    add_card(slide6, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "⛓️ Settlement & Agents", 
             ["• Algorand Testnet (AVM)",
              "• Testnet USDC ASA 10458941",
              "• 4 Sub-Agent Services (FastAPI + n8n)",
              "• SHA-256 Attestation Engine"], 
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 7: Demo Screenshots & Transaction Proof
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide7)
    add_header(slide7, "On-Chain Transaction Proof & Live Verification")

    add_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "🔗 Live Algorand Testnet Settlement", 
             ["• Router Wallet: 4DTSNS35EP...JDEHUHUNFR3KJGPO4",
              "• Testnet USDC ASA ID: 10458941",
              "• 5-Txn Atomic Payout Group executed in 1 single block.",
              "• Lora Block Explorer Link: https://lora.algokit.io/testnet"], 
             CYAN_ACCENT)

    add_card(slide7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "🔒 GoPlausible Verification & Receipt", 
             ["• GoPlausible Facilitator API (https://facilitator.goplausible.xyz) POST /verify & /settle verified.",
              "• SHA-256 Attestation Digest: sha256(token:score:verdict:txId:timestamp).",
              "• Live API Endpoint: https://api.dhanrajgupta.xyz/api/v1/health"], 
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 8: Innovation & Differentiation
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide8)
    add_header(slide8, "Innovation & Competitive Differentiation")

    add_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "Traditional Signal Services", 
             ["❌ High $300/mo subscription barrier.",
              "❌ Pay upfront for potentially dead/failing APIs.",
              "❌ Single-model bias & hallucination risks.",
              "❌ Centralized black-box server execution."], 
             RGBColor(239, 68, 68))

    add_card(slide8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "QuantMesh x402 Solution", 
             ["✅ $0.007 pay-per-use micropayments.",
              "✅ Zero-Fee Guarantee: $0 charged if an agent fails.",
              "✅ Multi-agent consensus weighted scoring.",
              "✅ Decentralized 5-txn atomic settlement in 1 block."], 
             GREEN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 9: Completeness & Functionality
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide9)
    add_header(slide9, "Completeness Status & Working Deliverables")

    add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "✅ Completed & Working Live", 
             ["• End-to-end HTTP 402 Probe-Challenge-Sign-Retry flow.",
              "• All 4 sub-agent microservices online on AWS EC2 & n8n.",
              "• Lute Wallet integration with indexesToSign=[0].",
              "• GoPlausible Facilitator /verify & /settle integration.",
              "• SHA-256 Attestation Digest generation."], 
             GREEN_ACCENT)

    add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "🚀 Planned Next Steps", 
             ["• Deploy PyTeal ARC-4 Box Storage contract to Mainnet.",
              "• Add dynamic worker registration portal for 3rd-party devs.",
              "• Integrate Pyth / Chainlink live oracle feeds.",
              "• Expand token support to 50+ crypto pairs."], 
             CYAN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 10: Real-World Impact & Scalability
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide10)
    add_header(slide10, "Real-World Impact, Scalability & Mainnet Readiness")

    add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "🌐 Real-World Impact & Ecosystem", 
             ["• Democratizes high-frequency AI signal access for retail traders.",
              "• Powers autonomous AI agents with machine-to-machine micropayments.",
              "• Creates an open monetization marketplace for AI developers on Algorand."], 
             CYAN_ACCENT)

    add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
             "⚡ Mainnet Readiness & Scalability", 
             ["• Algorand Mainnet Switch: Single flag flip to algorand:mainnet.",
              "• Mainnet USDC ASA 315667040 pre-configured.",
              "• Sub-3-second block finality & <$0.001 network fees enable infinite throughput."], 
             PURPLE_ACCENT)

    output_path = os.path.join("docs", "AlgoVerse2026_QuantMesh_x402_Official_Pitch.pptx")
    prs.save(output_path)
    print(f"Successfully generated: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
